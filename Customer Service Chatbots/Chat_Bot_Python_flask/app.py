from flask import Flask, request, jsonify, render_template
import pandas as pd
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
#parrot paraphrase
app = Flask(__name__)

# File paths
EXCEL_FILE = 'D:/N D/Bot/data/data.xlsx'
PDF_FILE = 'D:/N D/Bot/data/data.pdf'
TXT_FILE = 'D:/N D/Bot/data/data.txt'
PPT_FILE = 'D:/N D/Bot/data/data.pptx'
WORD_FILE = 'D:/N D/Bot/data/data.docx'

# Load data from Excel
def load_excel_data():
    df = pd.read_excel(EXCEL_FILE)
    return df

# Extract text from PDF
def extract_text_from_pdf():
    text = ""
    with open(PDF_FILE, 'rb') as file:
        reader = PdfReader(file)
        for page in reader.pages:
            text += page.extract_text()
    return text

# Extract text from TXT
def extract_text_from_txt():
    with open(TXT_FILE, 'r') as file:
        text = file.read()
    return text

# Extract text from PPT
def extract_text_from_ppt():
    text = ""
    presentation = Presentation(PPT_FILE)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text
    return text

# Extract text from Word
def extract_text_from_word():
    text = ""
    doc = Document(WORD_FILE)
    for para in doc.paragraphs:
        text += para.text
    return text

# Get response from different data sources
def get_response(user_input):
    # Check Excel
    df = load_excel_data()
    row = df[df['Name'].str.lower() == user_input.lower()]
    if not row.empty:
        return row.iloc[0]['Info']
    
    # Check PDF
    pdf_text = extract_text_from_pdf()
    if user_input.lower() in pdf_text.lower():
        return "Found in PDF"

    # Check TXT
    txt_text = extract_text_from_txt()
    if user_input.lower() in txt_text.lower():
        return "Found in TXT"

    # Check PPT
    ppt_text = extract_text_from_ppt()
    if user_input.lower() in ppt_text.lower():
        return "Found in PPT"

    # Check Word
    word_text = extract_text_from_word()
    if user_input.lower() in word_text.lower():
        return "Found in Word"

    return "Sorry, I don't have information on that."

# Home route to serve the HTML page
@app.route('/')
def home():
    return render_template('index.html')

# Webhook endpoint for chatbot
@app.route('/chat', methods=['POST'])
def chat():
    data = request.json
    user_input = data['userInput']
    response = get_response(user_input)
    return jsonify({"response": response})

if __name__ == '__main__':
    app.run(port=5000, debug=True)
